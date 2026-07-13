"""Regenerate PPT sample with 実題材 placeholders + no 'sample' footer."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

OUT = "/Users/anicca/.claude/skills/earn-gig/artifacts/5121769/ppt_sample.pptx"

DEEP_BLUE = RGBColor(0x1F, 0x4E, 0x79)
ORANGE = RGBColor(0xF2, 0xA1, 0x4E)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height

FONT = "Hiragino Sans"

def add_header(slide, text):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.8))
    shp.fill.solid(); shp.fill.fore_color.rgb = DEEP_BLUE
    shp.line.fill.background()
    tf = shp.text_frame; tf.margin_left = Inches(0.5); tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]; p.text = text
    r = p.runs[0]; r.font.name = FONT; r.font.size = Pt(20); r.font.color.rgb = WHITE; r.font.bold = True

def add_footer(slide, page_num):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.1), W, Inches(0.4))
    shp.fill.solid(); shp.fill.fore_color.rgb = ORANGE
    shp.line.fill.background()
    tf = shp.text_frame; tf.margin_left = Inches(0.5); tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]; p.text = f"算数 — 6年生 ／ Page {page_num}"
    r = p.runs[0]; r.font.name = FONT; r.font.size = Pt(11); r.font.color.rgb = WHITE

def add_text(slide, text, x, y, w, h, *, size=18, bold=False, color=None, align=PP_ALIGN.LEFT, name=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    r = p.runs[0]; r.font.name = name; r.font.size = Pt(size); r.font.bold = bold
    if color: r.font.color.rgb = color
    return tb

def add_bullet(tf, text, *, size=14, indent=0):
    p = tf.add_paragraph(); p.text = text; p.level = indent
    p.font.name = FONT; p.font.size = Pt(size)

# ─────────────────────────────────────────────────────
# Slide 1 — 表紙 (= 実題材: 算数 / 比例と反比例)
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s, "授業スライド ／ 第 3 回")
# Title
add_text(s, "比例と反比例 — 日常事例から考える",
         Inches(0.7), Inches(2.0), Inches(12), Inches(1.4),
         size=42, bold=True, color=DEEP_BLUE)
# Subtitle
add_text(s, "中学受験算数 ／ 数の関係を見つける",
         Inches(0.7), Inches(3.4), Inches(12), Inches(0.6),
         size=20, color=GRAY)
# Meta
add_text(s, "授業者：山田 直子 ／ 日付：2026 年 7 月 5 日",
         Inches(0.7), Inches(5.8), Inches(12), Inches(0.5),
         size=14, color=GRAY)
add_footer(s, 1)

# ─────────────────────────────────────────────────────
# Slide 2 — 本文 A (= 箇条書き + 図)
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s, "本日のポイント")
# Left: bullet list
tb = slide_left = s.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(6.0), Inches(5.5))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "学習のねらい"
p.font.name = FONT; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = DEEP_BLUE
add_bullet(tf, "比例・反比例の式を立てて、 表 と グラフ を読みとる", size=16)
add_bullet(tf, "日常事例 (買い物 ／ 速さ ／ 仕事率) から関係を見つける", size=16)
add_bullet(tf, "x が増えたとき y がどう変わるか説明できる", size=16)

add_bullet(tf, " ", size=12)
add_bullet(tf, "キーワード", size=18); tf.paragraphs[-1].font.bold = True; tf.paragraphs[-1].font.color.rgb = DEEP_BLUE
add_bullet(tf, "・比例定数  ・反比例  ・グラフ  ・表", size=14)

add_bullet(tf, " ", size=12)
add_bullet(tf, "授業の流れ", size=18); tf.paragraphs[-1].font.bold = True; tf.paragraphs[-1].font.color.rgb = DEEP_BLUE
add_bullet(tf, "① 導入：身近な事例で「変わり方」を観察", size=14)
add_bullet(tf, "② 展開：式 → 表 → グラフ の3表現を行き来", size=14)
add_bullet(tf, "③ 演習：3 問 → ペアで答え合わせ", size=14)
add_bullet(tf, "④ まとめ：今日の発見を 1 行で書く", size=14)

# Right: figure placeholder
fig = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.2), Inches(1.4), Inches(5.5), Inches(4.0))
fig.fill.solid(); fig.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
fig.line.color.rgb = GRAY
tf = fig.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "[ 図・画像 ]"
p.alignment = PP_ALIGN.CENTER; p.font.name = FONT; p.font.size = Pt(20); p.font.color.rgb = GRAY
p2 = tf.add_paragraph(); p2.text = "概念図 ／ 表 ／ グラフ をここに配置"
p2.alignment = PP_ALIGN.CENTER; p2.font.name = FONT; p2.font.size = Pt(14); p2.font.color.rgb = GRAY

# Right under fig: caption
add_text(s, "例：1 個 80 円のあめを x 個買うと、 代金 y 円 → y = 80x",
         Inches(7.2), Inches(5.6), Inches(5.5), Inches(0.5),
         size=12, color=GRAY)
# Right: confirm question box
qb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(6.0), Inches(5.5), Inches(0.9))
qb.fill.solid(); qb.fill.fore_color.rgb = ORANGE
qb.line.fill.background()
tf = qb.text_frame; tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.1)
p = tf.paragraphs[0]; p.text = "確認問題：x が 2 倍になると y はどうなる?"
p.font.name = FONT; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = WHITE
add_footer(s, 2)

# ─────────────────────────────────────────────────────
# Slide 3 — 本文 B (= 2 列 比較)
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s, "比較 — 比例 と 反比例")

# Left column
left = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.4), Inches(5.9), Inches(5.4))
left.fill.solid(); left.fill.fore_color.rgb = RGBColor(0xE9, 0xF1, 0xFA)
left.line.color.rgb = DEEP_BLUE
tf = left.text_frame; tf.margin_left = Inches(0.4); tf.margin_top = Inches(0.4)
p = tf.paragraphs[0]; p.text = "Ⅰ. 比例"
p.font.name = FONT; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = DEEP_BLUE
add_bullet(tf, " ", size=10)
add_bullet(tf, "・式： y = ax (a は比例定数)", size=16)
add_bullet(tf, "・x が 2 倍 → y も 2 倍", size=16)
add_bullet(tf, "・グラフ：原点を通る直線", size=16)
add_bullet(tf, " ", size=10)
add_bullet(tf, "例", size=18); tf.paragraphs[-1].font.bold = True; tf.paragraphs[-1].font.color.rgb = DEEP_BLUE
add_bullet(tf, "・あめの代金 (1 個 80 円)", size=14)
add_bullet(tf, "・移動距離 (一定の速さ)", size=14)

# Right column
right = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.4), Inches(5.9), Inches(5.4))
right.fill.solid(); right.fill.fore_color.rgb = RGBColor(0xFD, 0xEB, 0xD3)
right.line.color.rgb = ORANGE
tf = right.text_frame; tf.margin_left = Inches(0.4); tf.margin_top = Inches(0.4)
p = tf.paragraphs[0]; p.text = "Ⅱ. 反比例"
p.font.name = FONT; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = ORANGE
add_bullet(tf, " ", size=10)
add_bullet(tf, "・式： y = a / x (a は反比例定数)", size=16)
add_bullet(tf, "・x が 2 倍 → y は 1/2", size=16)
add_bullet(tf, "・グラフ：双曲線", size=16)
add_bullet(tf, " ", size=10)
add_bullet(tf, "例", size=18); tf.paragraphs[-1].font.bold = True; tf.paragraphs[-1].font.color.rgb = ORANGE
add_bullet(tf, "・1 つの仕事を何人かで分担", size=14)
add_bullet(tf, "・一定距離を歩く時間 (速さ vs 時間)", size=14)
add_footer(s, 3)

prs.save(OUT)
import os
print(f"saved: {OUT} ({os.path.getsize(OUT)} bytes)")
