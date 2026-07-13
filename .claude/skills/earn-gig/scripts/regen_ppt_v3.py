"""PPTX v3 — fix the 3 defects I found by verifying my own v2 output:
   (1) slide3 light-on-light text = unreadable -> solid dark headers + dark body on white
   (2) slide2 huge empty gray placeholder -> a REAL mini proportional graph drawn with shapes
   (3) plainness -> side accent band, real mini diagrams (比例 line / 反比例 curve), stronger hierarchy.
Education/授業 template, master-style, fully editable placeholders.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = "/Users/anicca/.claude/skills/earn-gig/artifacts/5121769/ppt_sample.pptx"
DEEP = RGBColor(0x1F, 0x4E, 0x79)
DEEP2 = RGBColor(0x2E, 0x6C, 0xA4)
ORANGE = RGBColor(0xE8, 0x82, 0x2B)
LIGHTBLUE = RGBColor(0xEA, 0xF1, 0xF8)
LIGHTORANGE = RGBColor(0xFD, 0xF0, 0xE2)
INK = RGBColor(0x22, 0x2A, 0x33)
GRAY = RGBColor(0x6B, 0x76, 0x82)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Hiragino Sans"

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def rect(slide, x, y, w, h, color, line=None, round=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp

def text(slide, s, x, y, w, h, *, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=None):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    first = True
    for line in s.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.text = line; p.alignment = align
        if sp is not None: p.space_after = Pt(sp)
        for r in p.runs:
            r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tb

def header(slide, title, page):
    rect(slide, 0, 0, W, Inches(0.95), DEEP)
    rect(slide, 0, Inches(0.95), W, Inches(0.06), ORANGE)
    text(slide, title, Inches(0.55), Inches(0.16), Inches(11), Inches(0.6), size=22, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, f"算数 6年 ／ {page}", Inches(11.0), Inches(7.06), Inches(2.2), Inches(0.35), size=10, color=GRAY, align=PP_ALIGN.RIGHT)

def mini_proportional_graph(slide, x, y, w, h, kind="line"):
    """Draw a real mini graph (axes + proportional line or inverse curve) with shapes."""
    rect(slide, x, y, w, h, WHITE, line=RGBColor(0xD0, 0xD7, 0xDE), round=True)
    pad = Inches(0.28)
    ax0x, ax0y = x + pad, y + h - pad
    # y axis
    yax = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, ax0x, y + pad, ax0x, ax0y)
    yax.line.color.rgb = GRAY; yax.line.width = Pt(1.25)
    # x axis
    xax = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, ax0x, ax0y, x + w - pad, ax0y)
    xax.line.color.rgb = GRAY; xax.line.width = Pt(1.25)
    if kind == "line":
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, ax0x, ax0y, x + w - pad, y + pad)
        ln.line.color.rgb = DEEP; ln.line.width = Pt(2.5)
    else:
        # inverse curve approximated by a freeform-ish elbow using a smooth arc shape
        arc = slide.shapes.add_shape(MSO_SHAPE.ARC, ax0x, y + pad - Emu(int(h*0.15)), int((w-2*pad)*1.7), int((h-2*pad)*1.7))
        arc.fill.background(); arc.line.color.rgb = ORANGE; arc.line.width = Pt(2.5); arc.shadow.inherit = False

# ── Slide 1: 表紙 (side accent band + title) ──
s1 = prs.slides.add_slide(BLANK)
rect(s1, 0, 0, Inches(0.45), H, DEEP)               # left accent band
rect(s1, Inches(0.45), 0, Inches(0.12), H, ORANGE)
text(s1, "授業スライド ／ 第 3 回", Inches(1.0), Inches(1.7), Inches(10), Inches(0.5), size=16, color=ORANGE, bold=True)
text(s1, "比例と反比例 — 日常事例から考える", Inches(1.0), Inches(2.3), Inches(11.5), Inches(1.6), size=40, bold=True, color=DEEP)
text(s1, "中学受験算数 ／ 数の関係を見つける", Inches(1.0), Inches(3.9), Inches(10), Inches(0.6), size=20, color=GRAY)
rect(s1, Inches(1.0), Inches(4.8), Inches(3.4), Inches(0.04), RGBColor(0xD0,0xD7,0xDE))
text(s1, "授業者：山田 直子　／　日付：2026 年 7 月 5 日", Inches(1.0), Inches(5.1), Inches(10), Inches(0.5), size=14, color=GRAY)
# decorative mini graph on cover (visual, not empty)
mini_proportional_graph(s1, Inches(9.7), Inches(4.7), Inches(2.6), Inches(2.1), "line")

# ── Slide 2: 本文A (ねらい/流れ + REAL graph, no empty box) ──
s2 = prs.slides.add_slide(BLANK); header(s2, "本日のポイント", "Page 2")
text(s2, "学習のねらい", Inches(0.55), Inches(1.3), Inches(6), Inches(0.4), size=18, bold=True, color=DEEP)
text(s2, "比例・反比例の式を立てて、表とグラフを読みとる\n日常事例（買い物 ／ 速さ ／ 仕事率）から関係を見つける\nx が増えたとき y がどう変わるか説明できる",
     Inches(0.6), Inches(1.75), Inches(6.4), Inches(1.6), size=15, color=INK, sp=4)
text(s2, "キーワード", Inches(0.55), Inches(3.35), Inches(6), Inches(0.4), size=18, bold=True, color=DEEP)
text(s2, "比例定数　・　反比例　・　グラフ　・　表", Inches(0.6), Inches(3.8), Inches(6.4), Inches(0.5), size=15, color=INK)
text(s2, "授業の流れ", Inches(0.55), Inches(4.45), Inches(6), Inches(0.4), size=18, bold=True, color=DEEP)
for n, t in enumerate(["導入：身近な事例で「変わり方」を観察", "展開：式 → 表 → グラフ の3表現を行き来",
                       "演習：3 問 → ペアで答え合わせ", "まとめ：今日の発見を 1 行で書く"]):
    rect(s2, Inches(0.6), Inches(4.95 + n*0.45), Inches(0.32), Inches(0.32), DEEP, round=True)
    text(s2, str(n+1), Inches(0.6), Inches(4.93 + n*0.45), Inches(0.32), Inches(0.34), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s2, t, Inches(1.05), Inches(4.95 + n*0.45), Inches(6), Inches(0.36), size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)
# right: a REAL proportional graph (replaces empty placeholder)
text(s2, "グラフ：y = 80x（あめ 1 個 80 円）", Inches(7.4), Inches(1.4), Inches(5.4), Inches(0.4), size=14, bold=True, color=GRAY)
mini_proportional_graph(s2, Inches(7.4), Inches(1.85), Inches(5.4), Inches(3.4), "line")
q = rect(s2, Inches(7.4), Inches(5.5), Inches(5.4), Inches(0.85), ORANGE, round=True)
text(s2, "確認問題：x が 2 倍になると y はどうなる？", Inches(7.5), Inches(5.5), Inches(5.2), Inches(0.85), size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ── Slide 3: 本文B (2列比較, FIXED contrast: solid dark headers + dark body on white) ──
s3 = prs.slides.add_slide(BLANK); header(s3, "比較 — 比例 と 反比例", "Page 3")
def compare_card(x, accent, light, head, eq, ex, kind):
    rect(s3, x, Inches(1.45), Inches(5.9), Inches(5.4), light, line=RGBColor(0xD8,0xDF,0xE6), round=True)
    rect(s3, x, Inches(1.45), Inches(5.9), Inches(0.85), accent, round=True)        # solid header bar
    text(s3, head, x, Inches(1.45), Inches(5.9), Inches(0.85), size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s3, "式", x+Inches(0.35), Inches(2.55), Inches(5.2), Inches(0.4), size=15, bold=True, color=accent)
    text(s3, eq, x+Inches(0.35), Inches(2.95), Inches(5.2), Inches(0.5), size=17, bold=True, color=INK)
    text(s3, "例", x+Inches(0.35), Inches(3.6), Inches(5.2), Inches(0.4), size=15, bold=True, color=accent)
    text(s3, ex, x+Inches(0.35), Inches(4.0), Inches(5.2), Inches(0.9), size=14, color=INK, sp=3)
    mini_proportional_graph(s3, x+Inches(1.7), Inches(5.0), Inches(2.5), Inches(1.65), kind)
compare_card(Inches(0.5), DEEP, LIGHTBLUE, "Ⅰ. 比例", "y = a x （a は比例定数）", "あめ 1 個 80 円 → y = 80x\nx（個）が増えると y（円）も増える", "line")
compare_card(Inches(6.9), ORANGE, LIGHTORANGE, "Ⅱ. 反比例", "y = a ÷ x （a は一定）", "12km を時速 x で進む → y = 12 ÷ x\nx（速さ）が増えると y（時間）は減る", "curve")

prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
